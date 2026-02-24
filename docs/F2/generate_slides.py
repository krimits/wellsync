"""
WellSync — Φ2 Slide Generator
Δημιουργεί το wellsync-phi2-presentation.pptx (11 slides, 16:9)
με τα σωστά χρώματα, κείμενα και layouts.

Run:  python docs/F2/generate_slides.py
Output: docs/F2/wellsync-phi2-presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Output path ──────────────────────────────────────────────────────────────
OUTPUT = "docs/F2/wellsync-phi2-presentation.pptx"

# ── Color palette ─────────────────────────────────────────────────────────────
PRIMARY   = RGBColor(0x1B, 0x4D, 0x4A)   # dark teal-green  #1B4D4A
ACCENT    = RGBColor(0x7B, 0xC6, 0x7E)   # lime-green       #7BC67E
BG        = RGBColor(0xF5, 0xF5, 0xF0)   # off-white        #F5F5F0
TEXT      = RGBColor(0x1A, 0x1A, 0x1A)   # near-black       #1A1A1A
TEXT_SEC  = RGBColor(0x6B, 0x72, 0x80)   # grey             #6B7280
AMBER     = RGBColor(0xF5, 0x9E, 0x0B)   # amber/gold       #F59E0B
AMBER_BG  = RGBColor(0xFE, 0xF3, 0xC7)   # pale amber       #FEF3C7
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
RED       = RGBColor(0xEF, 0x44, 0x44)   # error red        #EF4444
LIGHT_BG  = RGBColor(0xF0, 0xF0, 0xEB)   # light grey-green #F0F0EB

# ── Slide dimensions: 16:9 widescreen ────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

# ── Helper functions ──────────────────────────────────────────────────────────

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs: Presentation):
    """Add a completely blank slide."""
    blank_layout = prs.slide_layouts[6]   # index 6 = blank
    return prs.slides.add_slide(blank_layout)


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width_pt=0):
    """Add a filled/outlined rectangle shape."""
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, italic=False, color=TEXT,
                 align=PP_ALIGN.LEFT, wrap=True, font_name="Calibri"):
    """Add a text box with a single paragraph."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_slide_title(slide, title_text, y_offset=Inches(0.3)):
    """Add a prominent slide title."""
    add_text_box(
        slide, title_text,
        left=Inches(0.6), top=y_offset,
        width=Inches(12), height=Inches(0.7),
        font_size=28, bold=True, color=PRIMARY,
        align=PP_ALIGN.LEFT
    )


def add_accent_bar(slide):
    """Add a small accent color bar under the title."""
    add_rect(slide,
             left=Inches(0.6), top=Inches(0.95),
             width=Inches(1.2), height=Inches(0.06),
             fill_color=ACCENT)


def set_slide_bg(slide, color: RGBColor):
    """Set slide background to a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_bullet_box(slide, bullets: list[tuple[str, bool]],
                   left, top, width, height,
                   font_size=16, base_color=TEXT, indent_color=TEXT_SEC):
    """
    Add a text box with multiple bullet lines.
    bullets: list of (text, is_bold) tuples.
    Lines starting with '  ' are treated as sub-bullets.
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    for text, bold in bullets:
        is_sub = text.startswith("  ")
        clean = text.lstrip()
        prefix = "    • " if is_sub else "• "
        color = indent_color if is_sub else base_color
        sz = font_size - 2 if is_sub else font_size

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        run = p.add_run()
        run.text = prefix + clean
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return txBox


# ── SLIDE 1 — Title ───────────────────────────────────────────────────────────

def slide_01_title(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, PRIMARY)

    # Accent stripe at bottom
    add_rect(slide, 0, H - Inches(0.5), W, Inches(0.5), fill_color=ACCENT)

    # Main title
    add_text_box(
        slide, "WellSync",
        left=Inches(1.5), top=Inches(1.8),
        width=Inches(10), height=Inches(1.5),
        font_size=72, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER
    )

    # Tagline
    add_text_box(
        slide, "AI Holistic Wellness Coach",
        left=Inches(1.5), top=Inches(3.2),
        width=Inches(10), height=Inches(0.8),
        font_size=32, bold=False, italic=True, color=ACCENT,
        align=PP_ALIGN.CENTER
    )

    # Subtitle
    add_text_box(
        slide,
        "Ανάπτυξη Εφαρμογών Πληροφοριακών Συστημάτων  |  AUEB — Εαρινό 2026",
        left=Inches(1.5), top=Inches(4.3),
        width=Inches(10), height=Inches(0.5),
        font_size=16, color=RGBColor(0xCC, 0xCC, 0xCC),
        align=PP_ALIGN.CENTER
    )

    add_text_box(
        slide, "[Ονόματα ομάδας]",
        left=Inches(1.5), top=Inches(4.9),
        width=Inches(10), height=Inches(0.5),
        font_size=18, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER
    )


# ── SLIDE 2 — Το Πρόβλημα ────────────────────────────────────────────────────

def slide_02_problem(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, BG)
    add_slide_title(slide, "Γνωρίζεις πώς σε επηρεάζει ο ύπνος στην προπόνηση;")
    add_accent_bar(slide)

    # Three app "silos" boxes
    apps = [
        ("🍎 MyFitnessPal", "Διατροφή μόνο"),
        ("🚴 Strava", "Workout μόνο"),
        ("🧘 Headspace", "Mental health μόνο"),
    ]
    x_positions = [Inches(0.6), Inches(4.6), Inches(8.6)]
    for (name, desc), x in zip(apps, x_positions):
        add_rect(slide, x, Inches(1.3), Inches(3.6), Inches(2.2),
                 fill_color=WHITE, line_color=RGBColor(0xCC, 0xCC, 0xCC), line_width_pt=1)
        add_text_box(slide, name, x + Inches(0.15), Inches(1.5),
                     Inches(3.3), Inches(0.6), font_size=18, bold=True, color=PRIMARY)
        add_text_box(slide, desc, x + Inches(0.15), Inches(2.0),
                     Inches(3.3), Inches(0.5), font_size=14, color=TEXT_SEC)

    # "No connection" label between boxes
    add_text_box(slide, "✗ Καμία σύνδεση", Inches(0.6), Inches(3.7),
                 Inches(12), Inches(0.5), font_size=22, bold=True,
                 color=RED, align=PP_ALIGN.CENTER)

    # Key bullet points
    add_bullet_box(slide, [
        ("Κοιμάσαι 5h → πας γυμναστήριο → αποτυγχάνεις → δεν ξέρεις γιατί", False),
        ("3 εφαρμογές ανοιχτές — καμία δεν τις συνδέει", False),
        ("Καμία δεν σε ξέρει πραγματικά", True),
    ], Inches(0.6), Inches(4.4), Inches(12), Inches(2.5), font_size=17)


# ── SLIDE 3 — Η Λύση ─────────────────────────────────────────────────────────

def slide_03_solution(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, BG)
    add_slide_title(slide, "WellSync — Μία εφαρμογή που σε καταλαβαίνει")
    add_accent_bar(slide)

    # Left: bullets
    add_bullet_box(slide, [
        ("Καθημερινό 2-λεπτο check-in: ύπνος · διάθεση · ενέργεια · στρες", False),
        ("Καταγραφή προπόνησης και γευμάτων", False),
        ("AI Readiness Score (0–100) κάθε πρωί", True),
        ("Εξατομικευμένη σύσταση σε φυσική γλώσσα — όχι template", True),
    ], Inches(0.6), Inches(1.2), Inches(5.5), Inches(2.5), font_size=16)

    # Right: AI quote card
    add_rect(slide, Inches(6.4), Inches(1.1), Inches(6.5), Inches(4.2),
             fill_color=AMBER_BG, line_color=AMBER, line_width_pt=2)

    add_text_box(slide, "✨  AI Σύσταση (06:14)",
                 Inches(6.7), Inches(1.2), Inches(6.0), Inches(0.5),
                 font_size=13, bold=True, color=AMBER)

    quote = (
        '"Κοιμήθηκες 5.5h και χθες ήταν\n'
        'high intensity.\n\n'
        'Σήμερα: yoga 20\', πρωτεΐνη στο γεύμα.\n\n'
        'Αποφύγαμε HIIT — επέστρεψε στη\n'
        'ρουτίνα σου αύριο."'
    )
    add_text_box(slide, quote,
                 Inches(6.7), Inches(1.7), Inches(6.0), Inches(3.4),
                 font_size=17, italic=True, color=TEXT, align=PP_ALIGN.LEFT)

    add_text_box(slide, "↑ Παράχθηκε για εσένα, σήμερα — δεν είναι template",
                 Inches(0.6), Inches(5.5), Inches(12), Inches(0.5),
                 font_size=14, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)


# ── SLIDE 4 — User Flow ───────────────────────────────────────────────────────

def slide_04_flow(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, BG)
    add_slide_title(slide, "Η καθημερινή ροή")
    add_accent_bar(slide)

    # Three time-of-day columns
    cols = [
        ("🌅 Πρωί", ["Check-in (2λ)", "5 sliders", "→ Readiness Score", "→ AI Σύσταση"]),
        ("☀️ Μέρα",  ["Log Workout", "Log Meal", "< 1 λεπτό/καταγραφή"]),
        ("🌙 Βράδυ", ["Insights screen", "Trend graphs", '"Ύπνος > 7h → +23% energy"']),
    ]
    x_positions = [Inches(0.5), Inches(4.6), Inches(8.7)]
    for (header, items), x in zip(cols, x_positions):
        add_rect(slide, x, Inches(1.3), Inches(3.8), Inches(4.5),
                 fill_color=LIGHT_BG, line_color=ACCENT, line_width_pt=1.5)
        add_text_box(slide, header, x + Inches(0.15), Inches(1.45),
                     Inches(3.5), Inches(0.6), font_size=20, bold=True, color=PRIMARY)
        bullets = [(item, False) for item in items]
        add_bullet_box(slide, bullets, x + Inches(0.15), Inches(2.1),
                       Inches(3.5), Inches(3.4), font_size=15)

    # Arrow between columns
    for x_arr in [Inches(4.25), Inches(8.35)]:
        add_text_box(slide, "→", x_arr, Inches(3.2), Inches(0.5), Inches(0.6),
                     font_size=28, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ── SLIDE 5 — ML Innovation ───────────────────────────────────────────────────

def slide_05_ml(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, BG)
    add_slide_title(slide, "ML που συνδέει τομείς — per-user, όχι generic")
    add_accent_bar(slide)

    # Left: bullets
    add_bullet_box(slide, [
        ("Ridge Regression ανά χρήστη (scikit-learn)", False),
        ("Ανακαλύπτει ατομικά patterns:", True),
        ("  Πόσο επηρεάζει ο δικός σου ύπνος την απόδοσή σου;", False),
        ("  Πότε είσαι εσύ πιο productive στο gym;", False),
        ("Cold start: rule-based formula (αξία από μέρα 1)", False),
        ("Personalized mode: ενεργοποιείται μετά 14 ημέρες", True),
    ], Inches(0.6), Inches(1.2), Inches(6.0), Inches(4.0), font_size=16)

    # Right: example insight card
    add_rect(slide, Inches(7.0), Inches(1.2), Inches(5.8), Inches(3.0),
             fill_color=PRIMARY, line_color=None)
    add_text_box(slide, "📊 Insight (για εσένα)",
                 Inches(7.2), Inches(1.35), Inches(5.4), Inches(0.5),
                 font_size=14, bold=True, color=ACCENT)
    add_text_box(slide,
                 '"Ύπνος < 6h\n→ performance drop 38%\n\n(Δικό σου μοντέλο, 21 samples)"',
                 Inches(7.2), Inches(1.85), Inches(5.4), Inches(1.8),
                 font_size=19, italic=True, color=WHITE)

    # Cold start vs personalized visual
    add_rect(slide, Inches(0.6), Inches(5.4), Inches(5.0), Inches(0.7),
             fill_color=LIGHT_BG, line_color=ACCENT, line_width_pt=1)
    add_text_box(slide, "Ημέρες 1–13: Rule-based  →  Ημέρες 14+: ✅ Personalized ML",
                 Inches(0.75), Inches(5.5), Inches(4.8), Inches(0.5),
                 font_size=13, color=TEXT)


# ── SLIDE 6 — Agentic Loop ───────────────────────────────────────────────────

def slide_06_agents(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, BG)
    add_slide_title(slide, "Proactive AI — όχι passive app")
    add_accent_bar(slide)

    # Left: flowchart boxes
    flow_items = [
        (Inches(0.6),  Inches(1.2),  PRIMARY, WHITE,  "⏰  Scheduler 06:00"),
        (Inches(0.6),  Inches(2.05), PRIMARY, WHITE,  "📬  Event Queue"),
        (Inches(0.6),  Inches(2.9),  PRIMARY, WHITE,  "🔀  Event Gateway"),
        (Inches(0.6),  Inches(3.75), ACCENT,  TEXT,   "🤖  Morning Agent"),
        (Inches(0.6),  Inches(4.6),  AMBER,   TEXT,   "✨  Claude AI (enriched prompt)"),
        (Inches(0.6),  Inches(5.45), RGBColor(0x10,0x85,0x61), WHITE, "💾  agent_outputs → Frontend"),
    ]
    for x, y, fill, text_color, label in flow_items:
        add_rect(slide, x, y, Inches(5.2), Inches(0.65), fill_color=fill)
        add_text_box(slide, label, x + Inches(0.15), y + Inches(0.1),
                     Inches(4.9), Inches(0.5), font_size=16, bold=True, color=text_color)

    # Arrows
    for y_arr in [Inches(1.88), Inches(2.73), Inches(3.58), Inches(4.43), Inches(5.28)]:
        add_text_box(slide, "↓", Inches(1.7), y_arr, Inches(0.5), Inches(0.3),
                     font_size=16, color=TEXT_SEC, align=PP_ALIGN.CENTER)

    # Sub-items inside Morning Agent box
    add_text_box(slide,
                 "ML score  ·  RAG top-2 guidelines  ·  Prompt builder",
                 Inches(0.75), Inches(4.28), Inches(4.9), Inches(0.4),
                 font_size=11, color=TEXT, italic=True)

    # Right: bullets
    add_bullet_box(slide, [
        ("Το σύστημα δουλεύει ενώ εσύ κοιμάσαι", True),
        ("3 agents: πρωί · βράδυ · μεσάνυχτα (ML retraining)", False),
        ("RAG: ανακτά επιστημονικές οδηγίες (ACSM/WHO)", False),
        ("ML + RAG + Claude AI → evidence-based output", True),
        ("Fallback: αν το AI αποτύχει → ML-only (πάντα δουλεύει)", False),
    ], Inches(6.4), Inches(1.2), Inches(6.5), Inches(4.5), font_size=16)


# ── SLIDE 7 — Target Audience ─────────────────────────────────────────────────

def slide_07_audience(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, BG)
    add_slide_title(slide, "Για ποιον είναι το WellSync;")
    add_accent_bar(slide)

    # Left: Persona card
    add_rect(slide, Inches(0.6), Inches(1.2), Inches(5.8), Inches(5.0),
             fill_color=WHITE, line_color=PRIMARY, line_width_pt=2)
    add_text_box(slide, "👤  Νίκος, 21 ετών — Φοιτητής AUEB",
                 Inches(0.8), Inches(1.35), Inches(5.4), Inches(0.55),
                 font_size=17, bold=True, color=PRIMARY)

    persona_items = [
        ("✓  Πηγαίνει γυμναστήριο 2–3×/εβδ.", ACCENT),
        ("✓  Ασκείται ημι-τακτικά", ACCENT),
        ("✗  Ακανόνιστος ύπνος (εξεταστική)", RED),
        ("✗  Έχει εγκαταλείψει 3 wellness apps", RED),
        ("→  Θέλει εξήγηση, όχι απλά tracking", AMBER),
    ]
    for i, (text, color) in enumerate(persona_items):
        add_text_box(slide, text,
                     Inches(0.9), Inches(2.1) + Inches(0.55 * i),
                     Inches(5.2), Inches(0.5),
                     font_size=15, color=color)

    # Right: Market numbers
    add_text_box(slide, "Αγορά",
                 Inches(7.0), Inches(1.2), Inches(5.8), Inches(0.5),
                 font_size=22, bold=True, color=PRIMARY)

    market = [
        ("5.6B$", "Wellness app market (2024)"),
        ("15B$",  "Projected (2030)"),
        ("400K",  "Φοιτητές στην Ελλάδα"),
        ("1.8B",  "Millennials/Gen-Z με smartphone"),
    ]
    for i, (number, label) in enumerate(market):
        y = Inches(1.85) + Inches(1.1 * i)
        add_text_box(slide, number, Inches(7.0), y,
                     Inches(2.5), Inches(0.7),
                     font_size=36, bold=True, color=PRIMARY)
        add_text_box(slide, label, Inches(9.6), y + Inches(0.15),
                     Inches(3.2), Inches(0.5),
                     font_size=14, color=TEXT_SEC)


# ── SLIDE 8 — Competition ─────────────────────────────────────────────────────

def slide_08_competition(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, BG)
    add_slide_title(slide, "Γιατί όχι τα υπάρχοντα;")
    add_accent_bar(slide)

    headers = ["", "Workout", "Nutrition", "Mental", "Cross-domain\nML", "AI Rec."]
    rows = [
        ("MyFitnessPal", "±", "✓", "✗", "✗", "✗"),
        ("Strava",       "✓", "✗", "✗", "✗", "✗"),
        ("Headspace",    "✗", "✗", "✓", "✗", "✗"),
        ("Apple Health", "±", "±", "±", "✗", "✗"),
        ("WellSync ✅",  "✓", "✓", "✓", "✓", "✓"),
    ]

    col_widths = [Inches(2.2), Inches(1.7), Inches(1.7), Inches(1.7), Inches(2.2), Inches(1.7)]
    row_height = Inches(0.65)
    x_start = Inches(0.6)
    y_start = Inches(1.2)

    def cell_color(val, row_idx):
        if row_idx == 4:  # WellSync row
            return ACCENT
        return LIGHT_BG

    def text_color_for(val, row_idx):
        if row_idx == 4:
            return TEXT
        if val == "✓":
            return RGBColor(0x05, 0x96, 0x52)   # green
        if val == "✗":
            return RED
        return TEXT_SEC

    # Header row
    x = x_start
    for i, h in enumerate(headers):
        add_rect(slide, x, y_start, col_widths[i], row_height,
                 fill_color=PRIMARY)
        add_text_box(slide, h, x + Inches(0.05), y_start + Inches(0.1),
                     col_widths[i] - Inches(0.1), row_height - Inches(0.1),
                     font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x += col_widths[i]

    # Data rows
    for ri, row in enumerate(rows):
        x = x_start
        y = y_start + row_height * (ri + 1)
        bg = cell_color(None, ri)
        for ci, val in enumerate(row):
            add_rect(slide, x, y, col_widths[ci], row_height,
                     fill_color=bg,
                     line_color=RGBColor(0xCC, 0xCC, 0xCC), line_width_pt=0.5)
            tc = text_color_for(val, ri) if ci > 0 else (TEXT if ri < 4 else PRIMARY)
            add_text_box(slide, val, x + Inches(0.05), y + Inches(0.1),
                         col_widths[ci] - Inches(0.1), row_height - Inches(0.15),
                         font_size=14, bold=(ri == 4), color=tc,
                         align=(PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT))
            x += col_widths[ci]

    add_text_box(slide, "→ WellSync είναι το μόνο που συνδυάζει και τα 5",
                 Inches(0.6), Inches(6.5), Inches(12), Inches(0.5),
                 font_size=16, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)


# ── SLIDE 9 — Technology ──────────────────────────────────────────────────────

def slide_09_tech(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, BG)
    add_slide_title(slide, "Πώς το χτίζουμε")
    add_accent_bar(slide)

    # Left column: Backend/ML/RAG
    add_rect(slide, Inches(0.5), Inches(1.15), Inches(5.9), Inches(5.5),
             fill_color=WHITE, line_color=PRIMARY, line_width_pt=1.5)
    add_text_box(slide, "⚙️  Backend / ML / RAG",
                 Inches(0.7), Inches(1.25), Inches(5.5), Inches(0.5),
                 font_size=17, bold=True, color=PRIMARY)

    backend_items = [
        ("Python 3.10 + FastAPI", False),
        ("PostgreSQL + pgvector (vector search)", True),
        ("scikit-learn — Ridge Regression per user", False),
        ("sentence-transformers (semantic embeddings)", True),
        ("APScheduler + asyncio Queue", False),
        ("Claude AI API (Anthropic) — LLM συστάσεις", True),
    ]
    add_bullet_box(slide, backend_items,
                   Inches(0.7), Inches(1.85), Inches(5.5), Inches(4.0),
                   font_size=14)

    # Right column: Frontend/Deploy
    add_rect(slide, Inches(6.9), Inches(1.15), Inches(5.9), Inches(5.5),
             fill_color=WHITE, line_color=ACCENT, line_width_pt=1.5)
    add_text_box(slide, "🖥️  Frontend / Deploy",
                 Inches(7.1), Inches(1.25), Inches(5.5), Inches(0.5),
                 font_size=17, bold=True, color=PRIMARY)

    frontend_items = [
        ("React 18 + Vite (SPA)", False),
        ("Chart.js — trend graphs", False),
        ("JWT Authentication (python-jose)", False),
        ("Docker Compose (3 containers)", True),
        ("docker compose up --build → τρέχει", False),
    ]
    add_bullet_box(slide, frontend_items,
                   Inches(7.1), Inches(1.85), Inches(5.5), Inches(3.0),
                   font_size=14)

    # Architecture mini-diagram
    add_text_box(slide, "React  ↔  FastAPI  ↔  PostgreSQL+pgvector\n         ↕              ↕\n    ML Engine     Claude AI API",
                 Inches(7.1), Inches(5.0), Inches(5.5), Inches(1.5),
                 font_size=13, color=TEXT_SEC, italic=True)


# ── SLIDE 10 — Roadmap ────────────────────────────────────────────────────────

def slide_10_roadmap(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, BG)
    add_slide_title(slide, "Πού είμαστε — Πού πάμε")
    add_accent_bar(slide)

    phases = [
        ("Φεβρ 2026", "Ιδέα +\nΑρχιτεκτονική", True, "✓"),
        ("Μάρτιος 2026", "Φ1 Φόρμα ✓\nΦ2 Παρουσίαση ← ΤΩΡΑ", True, "▶"),
        ("Απρίλιος 2026", "Φ3: Technical\nDesign + Mockups", False, "○"),
        ("Μάιος 2026", "Φ4: Υλοποίηση\n+ Live Demo 🎯", False, "○"),
    ]

    y_line = Inches(3.2)
    x_positions = [Inches(1.2), Inches(4.4), Inches(7.6), Inches(10.8)]

    # Timeline line
    add_rect(slide, Inches(1.2), y_line + Inches(0.22), Inches(10.8), Inches(0.06),
             fill_color=RGBColor(0xCC, 0xCC, 0xCC))

    for i, ((month, desc, done, icon), x) in enumerate(zip(phases, x_positions)):
        fill = ACCENT if done else LIGHT_BG
        border = PRIMARY if done else RGBColor(0xCC, 0xCC, 0xCC)

        # Circle
        add_rect(slide, x - Inches(0.25), y_line, Inches(0.5), Inches(0.5),
                 fill_color=fill, line_color=border, line_width_pt=2)
        add_text_box(slide, icon, x - Inches(0.25), y_line + Inches(0.05),
                     Inches(0.5), Inches(0.4), font_size=14, bold=True,
                     color=WHITE if done else TEXT_SEC, align=PP_ALIGN.CENTER)

        # Month label above
        add_text_box(slide, month, x - Inches(1.1), y_line - Inches(0.7),
                     Inches(2.5), Inches(0.5), font_size=13, bold=True,
                     color=PRIMARY if done else TEXT_SEC, align=PP_ALIGN.CENTER)

        # Description below
        add_text_box(slide, desc, x - Inches(1.1), y_line + Inches(0.7),
                     Inches(2.5), Inches(1.0), font_size=13,
                     color=TEXT if done else TEXT_SEC, align=PP_ALIGN.CENTER)

    add_text_box(slide,
                 "Demo Φ4: live check-in → AI agent → εξατομικευμένη σύσταση — ζωντανά",
                 Inches(0.6), Inches(5.9), Inches(12), Inches(0.6),
                 font_size=17, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)


# ── SLIDE 11 — Closing ────────────────────────────────────────────────────────

def slide_11_closing(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, PRIMARY)

    # Accent bar at top
    add_rect(slide, 0, 0, W, Inches(0.4), fill_color=ACCENT)

    add_text_box(slide, "WellSync",
                 Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
                 font_size=64, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)

    add_text_box(slide,
                 '"Γνώρισε τον εαυτό σου — μέσα από τα δεδομένα σου."',
                 Inches(1.5), Inches(2.9), Inches(10), Inches(0.9),
                 font_size=24, italic=True, color=ACCENT,
                 align=PP_ALIGN.CENTER)

    add_text_box(slide, "[Ονόματα ομάδας]",
                 Inches(1.5), Inches(4.3), Inches(10), Inches(0.6),
                 font_size=20, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide, "Ερωτήσεις;",
                 Inches(1.5), Inches(5.3), Inches(10), Inches(0.8),
                 font_size=30, bold=True, color=AMBER,
                 align=PP_ALIGN.CENTER)

    # Accent bar at bottom
    add_rect(slide, 0, H - Inches(0.4), W, Inches(0.4), fill_color=ACCENT)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    print("Building slides...")
    slide_01_title(prs);         print("  [1/11] Title")
    slide_02_problem(prs);       print("  [2/11] Problem")
    slide_03_solution(prs);      print("  [3/11] Solution")
    slide_04_flow(prs);          print("  [4/11] User Flow")
    slide_05_ml(prs);            print("  [5/11] ML Innovation")
    slide_06_agents(prs);        print("  [6/11] Agentic Loop")
    slide_07_audience(prs);      print("  [7/11] Target Audience")
    slide_08_competition(prs);   print("  [8/11] Competition")
    slide_09_tech(prs);          print("  [9/11] Technology")
    slide_10_roadmap(prs);       print(" [10/11] Roadmap")
    slide_11_closing(prs);       print(" [11/11] Closing")

    prs.save(OUTPUT)
    print(f"\nSaved: {OUTPUT}")
    print("   -> Anoixe to se PowerPoint / Google Slides / Canva gia telikes veltiosis")


if __name__ == "__main__":
    main()
